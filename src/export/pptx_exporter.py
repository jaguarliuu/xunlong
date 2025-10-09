"""
PPTX - PPTPowerPoint
"""
import json
import re
from pathlib import Path
from typing import Dict, Any, Optional, List, Tuple
from loguru import logger
from html.parser import HTMLParser


class PPTXExporter:
    """PPTX"""

    async def export(
        self,
        project_dir: Path,
        output_path: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        PPTPPTX

        Args:
            project_dir: 
            output_path: 

        Returns:
            
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
            from pptx.dml.color import RGBColor

            # PPT
            ppt_data_file = project_dir / "reports" / "PPT_DATA.json"
            if not ppt_data_file.exists():
                return {
                    "status": "error",
                    "error": "PPT"
                }

            with open(ppt_data_file, 'r', encoding='utf-8') as f:
                ppt_data = json.load(f)

            logger.info(f"PPTX: {ppt_data.get('title', '')}")

            # PowerPoint
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # 
            colors = ppt_data.get('colors', {})
            primary_color = self._parse_color(colors.get('primary', '#3b82f6'))
            accent_color = self._parse_color(colors.get('accent', '#60a5fa'))
            text_color = self._parse_color(colors.get('text', '#1f2937'))
            bg_color = self._parse_color(colors.get('background', '#ffffff'))

            # 
            speech_notes = self._load_speech_notes(project_dir)

            # PPT
            slides_data = ppt_data.get('slides', [])
            for i, slide_data in enumerate(slides_data):
                slide_number = slide_data.get('slide_number', i + 1)
                html_content = slide_data.get('html_content', '')

                logger.info(f" {slide_number} ")

                # 
                blank_slide_layout = prs.slide_layouts[6]  # 
                slide = prs.slides.add_slide(blank_slide_layout)

                # 
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = bg_color

                # HTML
                self._render_slide_from_html(
                    slide,
                    html_content,
                    primary_color,
                    accent_color,
                    text_color
                )

                # 
                self._add_page_number(slide, slide_number, len(slides_data), text_color)

                # 
                if speech_notes and slide_number <= len(speech_notes):
                    notes_text = speech_notes[slide_number - 1].get('speech_notes', '')
                    if notes_text:
                        notes_slide = slide.notes_slide
                        text_frame = notes_slide.notes_text_frame
                        text_frame.text = notes_text

            # 
            if not output_path:
                output_path = project_dir / "exports" / f"{ppt_data.get('title', 'presentation')}.pptx"
            else:
                output_path = Path(output_path)

            # 
            output_path.parent.mkdir(parents=True, exist_ok=True)

            # 
            prs.save(str(output_path))

            # 
            file_size = output_path.stat().st_size
            file_size_str = self._format_file_size(file_size)

            logger.info(f"PPTX: {output_path}")

            return {
                "status": "success",
                "output_file": str(output_path),
                "file_size": file_size_str,
                "slide_count": len(slides_data)
            }

        except ImportError as e:
            return {
                "status": "error",
                "error": f"python-pptx: pip install python-pptx ({e})"
            }
        except Exception as e:
            logger.error(f"PPTX: {e}")
            import traceback
            traceback.print_exc()
            return {
                "status": "error",
                "error": str(e)
            }

    def _load_speech_notes(self, project_dir: Path) -> Optional[list]:
        """TODO: Add docstring."""
        speech_file = project_dir / "reports" / "SPEECH_NOTES.json"
        if speech_file.exists():
            with open(speech_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return data.get('speech_notes', [])
        return None

    def _render_slide_from_html(
        self,
        slide,
        html_content: str,
        primary_color,
        accent_color,
        text_color
    ):
        """
        HTML
        HTML
        """
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

        # HTML
        parser = HTMLStructureParser()
        parser.feed(html_content)

        elements = parser.elements

        if not elements:
            # 
            text = self._extract_plain_text(html_content)
            self._add_simple_text_box(slide, text, 0.5, 0.5, 9, 6.5, text_color, 18)
            return

        # 
        self._smart_layout(slide, elements, primary_color, accent_color, text_color)

    def _smart_layout(self, slide, elements, primary_color, accent_color, text_color):
        """TODO: Add docstring."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

        # 
        titles = [e for e in elements if e['type'] in ['h1', 'h2', 'h3']]
        content_items = [e for e in elements if e['type'] in ['p', 'li', 'div']]

        y_offset = 0.5  # Y

        # 
        if titles:
            title_elem = titles[0]
            title_text = title_elem['text']

            # 
            if title_elem['type'] == 'h1':
                font_size = 44
            elif title_elem['type'] == 'h2':
                font_size = 36
            else:
                font_size = 28

            # 
            title_box = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(y_offset),
                Inches(9),
                Inches(1.2)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True

            p = title_frame.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = primary_color
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            y_offset += 1.5

        # 
        if content_items:
            content_height = 7.5 - y_offset - 0.8  # 

            # 
            if len(content_items) <= 3:
                self._render_content_large(
                    slide,
                    content_items,
                    y_offset,
                    content_height,
                    text_color,
                    accent_color
                )
            elif len(content_items) <= 8:
                # 
                self._render_content_list(
                    slide,
                    content_items,
                    y_offset,
                    content_height,
                    text_color,
                    accent_color
                )
            else:
                # 
                self._render_content_compact(
                    slide,
                    content_items,
                    y_offset,
                    content_height,
                    text_color
                )

    def _render_content_large(self, slide, items, y_offset, height, text_color, accent_color):
        """TODO: Add docstring."""
        from pptx.util import Inches, Pt

        for i, item in enumerate(items[:5]):  # 5
            text_box = slide.shapes.add_textbox(
                Inches(1.0),
                Inches(y_offset + i * (height / 5)),
                Inches(8.0),
                Inches(height / 5 - 0.1)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = 1  # 

            p = text_frame.paragraphs[0]
            p.text = item['text']
            p.font.size = Pt(24)
            p.font.color.rgb = text_color

            # 
            if i < len(items):
                # 
                circle = slide.shapes.add_shape(
                    1,  # 
                    Inches(0.6),
                    Inches(y_offset + i * (height / 5) + 0.15),
                    Inches(0.2),
                    Inches(0.2)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = accent_color
                circle.line.color.rgb = accent_color

    def _render_content_list(self, slide, items, y_offset, height, text_color, accent_color):
        """TODO: Add docstring."""
        from pptx.util import Inches, Pt

        text_box = slide.shapes.add_textbox(
            Inches(1.0),
            Inches(y_offset),
            Inches(8.0),
            Inches(height)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(items[:12]):  # 12
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]

            p.text = f" {item['text']}"
            p.font.size = Pt(18)
            p.font.color.rgb = text_color
            p.space_after = Pt(8)
            p.level = 0

    def _render_content_compact(self, slide, items, y_offset, height, text_color):
        """TODO: Add docstring."""
        from pptx.util import Inches, Pt

        # 
        left_items = items[::2]  # 
        right_items = items[1::2]  # 

        # 
        left_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(y_offset),
            Inches(4.5),
            Inches(height)
        )
        self._fill_text_frame(left_box.text_frame, left_items, text_color, 14)

        # 
        right_box = slide.shapes.add_textbox(
            Inches(5.0),
            Inches(y_offset),
            Inches(4.5),
            Inches(height)
        )
        self._fill_text_frame(right_box.text_frame, right_items, text_color, 14)

    def _fill_text_frame(self, text_frame, items, text_color, font_size):
        """TODO: Add docstring."""
        from pptx.util import Pt

        text_frame.word_wrap = True

        for i, item in enumerate(items):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]

            p.text = f" {item['text']}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = text_color
            p.space_after = Pt(4)

    def _add_simple_text_box(self, slide, text, left, top, width, height, color, font_size):
        """TODO: Add docstring."""
        from pptx.util import Inches, Pt

        if not text:
            return

        text_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color

    def _add_page_number(self, slide, current, total, text_color):
        """TODO: Add docstring."""
        from pptx.util import Inches, Pt

        page_box = slide.shapes.add_textbox(
            Inches(8.5),
            Inches(7.0),
            Inches(1.0),
            Inches(0.3)
        )
        text_frame = page_box.text_frame

        p = text_frame.paragraphs[0]
        p.text = f"{current} / {total}"
        p.font.size = Pt(12)
        p.font.color.rgb = text_color

    def _extract_plain_text(self, html: str) -> str:
        """TODO: Add docstring."""
        # HTML
        text = re.sub(r'<[^>]+>', ' ', html)
        # 
        text = re.sub(r'\s+', ' ', text)
        return text.strip()

    def _parse_color(self, color_str: str):
        """RGB"""
        from pptx.dml.color import RGBColor

        if color_str.startswith('#'):
            color_str = color_str[1:]

        try:
            r = int(color_str[0:2], 16)
            g = int(color_str[2:4], 16)
            b = int(color_str[4:6], 16)
            return RGBColor(r, g, b)
        except:
            return RGBColor(59, 130, 246)  # 

    def _format_file_size(self, size_bytes: int) -> str:
        """TODO: Add docstring."""
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size_bytes < 1024.0:
                return f"{size_bytes:.2f} {unit}"
            size_bytes /= 1024.0
        return f"{size_bytes:.2f} TB"


class HTMLStructureParser(HTMLParser):
    """HTML"""

    def __init__(self):
        super().__init__()
        self.elements = []
        self.current_tag = None
        self.current_text = []

    def handle_starttag(self, tag, attrs):
        if tag in ['h1', 'h2', 'h3', 'p', 'li', 'div']:
            self.current_tag = tag
            self.current_text = []

    def handle_endtag(self, tag):
        if tag in ['h1', 'h2', 'h3', 'p', 'li', 'div'] and self.current_tag == tag:
            text = ''.join(self.current_text).strip()
            if text:
                self.elements.append({
                    'type': tag,
                    'text': text
                })
            self.current_tag = None
            self.current_text = []

    def handle_data(self, data):
        if self.current_tag:
            cleaned = data.strip()
            if cleaned:
                self.current_text.append(cleaned + ' ')
